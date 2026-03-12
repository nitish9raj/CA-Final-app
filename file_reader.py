"""
file_reader.py — Safe multi-format text extractor.
Handles PDF, DOCX, PPTX, and images (PNG/JPG).
Falls back gracefully if optional deps are missing.
"""
import io

def extract_text(uploaded_file):
    """Routes uploaded file to the correct parser. Returns extracted text string."""
    if uploaded_file is None:
        return ""
    filename = uploaded_file.name.lower()
    try:
        if filename.endswith(".pdf"):
            return _extract_pdf(uploaded_file)
        elif filename.endswith(".docx"):
            return _extract_docx(uploaded_file)
        elif filename.endswith(".pptx"):
            return _extract_pptx(uploaded_file)
        elif filename.endswith((".png", ".jpg", ".jpeg")):
            return _extract_image(uploaded_file)
        elif filename.endswith(".txt"):
            return uploaded_file.read().decode("utf-8", errors="ignore")
        else:
            return f"Unsupported file type: {filename}"
    except Exception as e:
        return f"Error extracting text: {e}"

def _extract_pdf(f):
    try:
        import pdfplumber
        text = ""
        with pdfplumber.open(f) as pdf:
            for page in pdf.pages:
                extracted = page.extract_text()
                if extracted:
                    text += extracted + "\n"
        return text.strip() or "No text found in PDF."
    except ImportError:
        # Fallback: try PyPDF2
        try:
            import PyPDF2
            reader = PyPDF2.PdfReader(f)
            return "\n".join(p.extract_text() or "" for p in reader.pages).strip()
        except ImportError:
            return "PDF parsing requires pdfplumber. Run: pip install pdfplumber"

def _extract_docx(f):
    try:
        import docx
        doc = docx.Document(f)
        return "\n".join(p.text for p in doc.paragraphs if p.text.strip())
    except ImportError:
        return "DOCX parsing requires python-docx. Run: pip install python-docx"

def _extract_pptx(f):
    try:
        from pptx import Presentation
        prs = Presentation(f)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text += shape.text + "\n"
        return text.strip() or "No text found in PPTX."
    except ImportError:
        return "PPTX parsing requires python-pptx. Run: pip install python-pptx"

def _extract_image(f):
    try:
        import pytesseract
        from PIL import Image
        image = Image.open(f)
        return pytesseract.image_to_string(image).strip() or "No text detected in image."
    except ImportError:
        return ("Image OCR requires pytesseract + Pillow.\n"
                "Run: pip install pytesseract Pillow\n"
                "Also install Tesseract OCR engine from: https://github.com/tesseract-ocr/tesseract")
    except Exception as e:
        return f"Image OCR error: {e}"

def get_word_count(text):
    return len(text.split()) if text else 0

def get_preview(text, chars=500):
    if not text: return ""
    return text[:chars] + ("..." if len(text) > chars else "")
